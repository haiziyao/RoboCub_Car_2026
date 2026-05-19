from pathlib import Path
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))

from build_midterm_ppt import (  # noqa: E402
    ASSET_DIR,
    BG,
    BG_2,
    BLUE,
    COPPER,
    FONT,
    GREEN,
    INK,
    LINE,
    MUTED,
    RED,
    SLIDE_H,
    SLIDE_W,
    TEAL,
    WHITE,
    add_bullets,
    add_chip,
    add_connector,
    add_footer,
    add_full_image,
    add_metric,
    add_node,
    add_overlay,
    add_panel,
    add_section_label,
    add_solid_block,
    add_table_header,
    add_table_row,
    add_textbox,
    add_title,
    rgb,
    set_bg,
)


OUT = ROOT / "RuboVision-Engine-midterm-report-v2.pptx"


def add_small_note(slide, text, x, y, w, color=MUTED):
    add_textbox(slide, text, x, y, w, Inches(0.24), 8.8, color)


def add_code_panel(slide, title, lines, x, y, w, h, accent=TEAL):
    add_panel(slide, x, y, w, h, fill=rgb("111722"), line=rgb("334155"))
    add_textbox(slide, title, x + Inches(0.22), y + Inches(0.18), w - Inches(0.44), Inches(0.22), 11, accent, True)
    text = "\n".join(lines)
    box = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.58), w - Inches(0.5), h - Inches(0.78))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.line_spacing = 1.08
    for run in p.runs:
        run.font.name = "Cascadia Mono"
        run.font.size = Pt(8.4)
        run.font.color.rgb = RGBColor(229, 234, 243)


def add_status_card(slide, title, value, note, x, y, w, accent):
    add_panel(slide, x, y, w, Inches(1.12), fill=WHITE)
    add_solid_block(slide, accent, x, y, w, Inches(0.08))
    add_textbox(slide, title, x + Inches(0.18), y + Inches(0.22), w - Inches(0.36), Inches(0.22), 10.2, MUTED, True)
    add_textbox(slide, value, x + Inches(0.18), y + Inches(0.48), w - Inches(0.36), Inches(0.26), 16, accent, True)
    add_textbox(slide, note, x + Inches(0.18), y + Inches(0.82), w - Inches(0.36), Inches(0.18), 8.2, MUTED)


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_image(slide, ASSET_DIR / "cover-vision-engine.png")
    add_overlay(slide, BG, transparency=10)
    add_solid_block(slide, rgb("0F172A"), 0, 0, Inches(5.95), SLIDE_H)
    add_solid_block(slide, TEAL, Inches(5.95), 0, Inches(0.035), SLIDE_H)
    add_chip(slide, "阶段更新版 V2", Inches(0.78), Inches(0.9), Inches(1.35), COPPER)
    add_textbox(slide, "RuboVision Engine", Inches(0.76), Inches(1.42), Inches(5.4), Inches(0.58), 30, WHITE, True)
    add_textbox(slide, "中期报告", Inches(0.78), Inches(2.1), Inches(3.2), Inches(0.46), 24, TEAL, True)
    add_textbox(
        slide,
        "配置驱动视觉任务框架 · Vision 模块重构 · Web 调试 · 串口联调验证",
        Inches(0.8),
        Inches(2.82),
        Inches(5.0),
        Inches(0.62),
        13,
        RGBColor(229, 235, 245),
    )
    add_bullets(
        slide,
        [
            "项目语言与框架：Rust + OpenCV + Axum",
            "报告时间：2026 年 5 月",
            "本版重点：结合最新重构要求、测试要求与联调问题复盘",
        ],
        Inches(0.82),
        Inches(4.85),
        Inches(5.05),
        Inches(1.2),
        11.2,
        RGBColor(231, 237, 247),
        gap=0.04,
        accent=COPPER,
    )
    add_footer(slide, 1, dark=True)


def slide_context(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_2)
    add_title(slide, "项目定位与阶段原则", "中期目标不只是让功能跑通，而是把视觉任务组织成清晰、可调、可联调的工程框架")
    add_metric(slide, "主链路", "5 层", "Source / Listener / Dispatcher / Executor / WebMessage", Inches(0.82), Inches(1.45), Inches(3.65), TEAL)
    add_metric(slide, "视觉模块", "6 文件", "camera / config / color / qr / cross / response", Inches(4.82), Inches(1.45), Inches(3.65), GREEN)
    add_metric(slide, "测试方向", "3 类", "单功能 GUI、Web 调试、虚拟串口联调", Inches(8.82), Inches(1.45), Inches(3.65), COPPER)
    add_section_label(slide, "阶段设计约束", Inches(0.88), Inches(3.35), COPPER)
    add_bullets(
        slide,
        [
            "参数类型由上层严格保证，因此功能实现应减少过度防御与泛用化包装。",
            "保留已有架构占位：WebSource、TaskExecutor、空 impl/default/TODO/dead_code 标记不随意删除。",
            "Vision 以可读、可测、可定位为优先，避免单文件膨胀导致找测试和找功能成本过高。",
            "测试函数必须接近配置文件入口：读取实际配置、直接 imshow、单一功能单独验证。",
        ],
        Inches(0.92),
        Inches(3.75),
        Inches(11.5),
        Inches(1.8),
        12.5,
        INK,
        gap=0.07,
        accent=TEAL,
    )
    add_footer(slide, 2)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_image(slide, ASSET_DIR / "framework-architecture.png")
    add_overlay(slide, BG, transparency=20)
    add_title(slide, "总体架构：事件进入与结果返回", "框架保留扩展位，但执行路径保持简单直接", dark=True)
    y = Inches(2.0)
    xs = [Inches(0.7), Inches(3.08), Inches(5.46), Inches(7.84), Inches(10.22)]
    labels = ["Source\nUART/Web/Timer", "TaskListener\n收事件", "Dispatcher\n找设备+函数", "TaskExecutor\n执行函数", "WebMessage\n返回结果"]
    colors = [TEAL, BLUE, COPPER, GREEN, RED]
    for i, (x, label, color) in enumerate(zip(xs, labels, colors)):
        add_node(slide, label, x, y, Inches(1.72), Inches(0.9), color, WHITE, 11.8)
        if i < len(xs) - 1:
            add_connector(slide, x + Inches(1.72), y + Inches(0.45), xs[i + 1], y + Inches(0.45), RGBColor(206, 224, 235))
    add_panel(slide, Inches(0.82), Inches(4.15), Inches(5.65), Inches(1.36), fill=rgb("111827"), line=rgb("3B4758"))
    add_textbox(slide, "配置输入", Inches(1.05), Inches(4.37), Inches(1.1), Inches(0.22), 11, COPPER, True)
    add_bullets(
        slide,
        [
            "bindings.toml：source_key → task/device/function",
            "device.toml：全局 UART 与 Camera path",
            "func_param.toml：函数参数与 returns.web/gpio",
        ],
        Inches(1.05),
        Inches(4.72),
        Inches(5.1),
        Inches(0.72),
        9.6,
        RGBColor(235, 240, 248),
        gap=0.01,
        accent=TEAL,
    )
    add_panel(slide, Inches(6.82), Inches(4.15), Inches(5.65), Inches(1.36), fill=rgb("111827"), line=rgb("3B4758"))
    add_textbox(slide, "输出通道", Inches(7.05), Inches(4.37), Inches(1.1), Inches(0.22), 11, TEAL, True)
    add_bullets(
        slide,
        [
            "Web：任务消息、调试图像、历史快照",
            "UART：识别完成后写回结果",
            "GPIO：树莓派环境下状态灯反馈",
        ],
        Inches(7.05),
        Inches(4.72),
        Inches(5.1),
        Inches(0.72),
        9.6,
        RGBColor(235, 240, 248),
        gap=0.01,
        accent=COPPER,
    )
    add_footer(slide, 3, dark=True)


def slide_returns_config(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_2)
    add_title(slide, "函数返回配置：结果通道与业务参数分离", "按照最新要求，web/gpio 是否返回作为 returns 独立配置，识别参数仍留在 args 中")
    add_code_panel(
        slide,
        "func_param.toml 示例",
        [
            "[[func_param_config.func_param_list]]",
            'function_id = "color_detect"',
            "returns = { web = true, gpio = true }",
            "args = [",
            '  "debug_model=false",',
            '  "loop_count=5",',
            '  "radius_ratio=0.4",',
            '  "color.red=0,50,160,255,110,255",',
            '  "color_light_pin=17",',
            "]",
        ],
        Inches(0.82),
        Inches(1.48),
        Inches(5.65),
        Inches(4.6),
        TEAL,
    )
    add_section_label(slide, "设计收益", Inches(7.0), Inches(1.55), COPPER)
    add_bullets(
        slide,
        [
            "语义更清楚：returns 只表示结果/反馈通道，不混入识别算法参数。",
            "测试更方便：单功能测试可关闭 web/gpio，仅验证视觉识别本身。",
            "联调更可控：实际运行时按部署环境启用 Web 或 GPIO 返回。",
            "实现更简洁：函数执行层只读取 returns.web / returns.gpio 做分支。",
        ],
        Inches(7.05),
        Inches(1.92),
        Inches(5.1),
        Inches(1.8),
        12,
        INK,
        gap=0.08,
        accent=GREEN,
    )
    add_status_card(slide, "代码锚点", "FuncReturnConfig", "src/config/func.rs", Inches(7.05), Inches(4.25), Inches(2.48), TEAL)
    add_status_card(slide, "执行锚点", "returns.web", "src/init/task_exec.rs", Inches(9.85), Inches(4.25), Inches(2.48), BLUE)
    add_status_card(slide, "硬件锚点", "returns.gpio", "src/func/usual.rs", Inches(8.45), Inches(5.55), Inches(2.48), COPPER)
    add_footer(slide, 4)


def slide_vision_refactor(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_title(slide, "Vision 模块重构：拆小文件，降低定位成本", "从“大文件堆叠”调整为按职责拆分，测试入口集中但功能实现分散", dark=True)
    files = [
        ("vision.rs", "模块出口\npub use 聚合", TEAL),
        ("camera.rs", "摄像头打开\nV4L2 入口", BLUE),
        ("config.rs", "Camera / Color / QR\n配置解析", COPPER),
        ("color.rs", "HSV 颜色识别\nROI + 稳定计数", GREEN),
        ("qr.rs", "二维码识别\n灰度 + quircs", RED),
        ("cross.rs", "黑色轮廓识别\n预留函数体", RGBColor(120, 132, 158)),
        ("response.rs", "UART/GPIO 回执\n串口写回", RGBColor(115, 184, 148)),
        ("tests.rs", "配置驱动测试\nimshow + Web 调试", RGBColor(95, 161, 217)),
    ]
    x0 = Inches(0.72)
    y0 = Inches(1.55)
    for i, (name, desc, color) in enumerate(files):
        x = x0 + Inches(3.12) * (i % 4)
        y = y0 + Inches(1.7) * (i // 4)
        add_panel(slide, x, y, Inches(2.78), Inches(1.18), fill=rgb("232936"), line=rgb("3A4657"))
        add_solid_block(slide, color, x, y, Inches(2.78), Inches(0.08))
        add_textbox(slide, name, x + Inches(0.18), y + Inches(0.25), Inches(2.3), Inches(0.24), 13, WHITE, True)
        add_textbox(slide, desc, x + Inches(0.18), y + Inches(0.58), Inches(2.3), Inches(0.42), 9.5, RGBColor(213, 221, 233))
    add_panel(slide, Inches(0.82), Inches(5.38), Inches(11.65), Inches(0.88), fill=rgb("111827"), line=rgb("3B4758"))
    add_textbox(
        slide,
        "重构边界：去掉冗余封装和过度泛用设计，但不删除项目既有架构占位；cross_detect 当前返回默认值，后续接入黑色轮廓识别。",
        Inches(1.08),
        Inches(5.65),
        Inches(11.0),
        Inches(0.25),
        11,
        RGBColor(235, 240, 248),
        True,
        PP_ALIGN.CENTER,
    )
    add_footer(slide, 5, dark=True)


def slide_vision_functions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_2)
    add_title(slide, "视觉功能状态：完整、可测、可继续扩展", "颜色与二维码按照参考项目完成主流程，cross 先保留稳定占位")
    cards = [
        ("颜色识别", "已完成", ["读取 color_camera", "圆形 ROI mask", "HSV 阈值由配置注册", "连续 stable count 后返回"], TEAL),
        ("二维码识别", "已完成", ["读取 qr_camera", "灰度仅用于解码", "显示窗口保持彩色帧", "payload 解析为任务编号"], GREEN),
        ("十字/路口识别", "占位", ["当前返回 0", "函数体保留", "后续识别黑色轮廓", "不影响编译与调度链路"], COPPER),
    ]
    for i, (title, status, bullets, color) in enumerate(cards):
        x = Inches(0.82) + Inches(4.18) * i
        add_panel(slide, x, Inches(1.55), Inches(3.78), Inches(4.72), fill=WHITE)
        add_solid_block(slide, color, x, Inches(1.55), Inches(3.78), Inches(0.1))
        add_textbox(slide, title, x + Inches(0.22), Inches(1.95), Inches(2.2), Inches(0.28), 16, INK, True)
        add_chip(slide, status, x + Inches(2.55), Inches(1.92), Inches(0.78), color)
        add_bullets(slide, bullets, x + Inches(0.28), Inches(2.72), Inches(3.1), Inches(1.7), 11.2, INK, gap=0.08, accent=color)
        add_small_note(slide, ["color.rs", "qr.rs", "cross.rs"][i], x + Inches(0.28), Inches(5.72), Inches(2.0), color)
    add_footer(slide, 6)


def slide_test_design(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_title(slide, "测试设计：单功能直测 + 配置驱动", "测试函数应像原项目一样直接打开窗口，不强制返回 Web/GPIO", dark=True)
    left = [
        "test_color_detect_with_config",
        "test_qr_detect_with_config",
        "test_cross_detect_with_config",
        "test_usb_camera_open_and_read_frame_from_config",
    ]
    add_section_label(slide, "单功能 GUI 测试", Inches(0.85), Inches(1.55), TEAL)
    add_bullets(slide, left, Inches(0.95), Inches(1.95), Inches(5.4), Inches(1.8), 12, RGBColor(235, 240, 248), gap=0.08, accent=TEAL)
    add_section_label(slide, "测试规则", Inches(6.95), Inches(1.55), COPPER)
    add_bullets(
        slide,
        [
            "读取 config/device.toml 与 config/func_param.toml。",
            "每个功能一个测试函数，避免在综合测试里混淆问题来源。",
            "使用 imshow 直观查看 draw frame 和最终识别状态。",
            "默认 ignored，需要摄像头/GUI 时由开发者显式运行。",
        ],
        Inches(7.05),
        Inches(1.95),
        Inches(5.2),
        Inches(1.8),
        12,
        RGBColor(235, 240, 248),
        gap=0.08,
        accent=COPPER,
    )
    add_code_panel(
        slide,
        "运行示例",
        [
            "cargo test --package new_car --lib \\",
            "  device::vision::tests::test_color_detect_with_config \\",
            "  -- --ignored --nocapture",
            "",
            "cargo test --package new_car --lib \\",
            "  device::vision::tests::test_qr_detect_with_config \\",
            "  -- --ignored --nocapture",
        ],
        Inches(1.0),
        Inches(4.45),
        Inches(11.35),
        Inches(1.35),
        GREEN,
    )
    add_footer(slide, 7, dark=True)


def slide_web_debug(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_2)
    add_title(slide, "Web 调试：复用原 Web 端传回 Vision 结果", "新增带 Web 的测试函数，不另造调试协议，使用原 /message 入口推送 WebMessage")
    add_code_panel(
        slide,
        "Vision Web 测试行为",
        [
            "1. 打开摄像头并持续识别",
            "2. 每次识别成功后绘制 draw frame",
            "3. 将最终帧编码为 JPEG base64 data URL",
            "4. WebMessage { text, image } POST 到 /message",
            "5. 发送成功后 sleep 5 秒，再继续识别",
        ],
        Inches(0.82),
        Inches(1.5),
        Inches(5.65),
        Inches(4.65),
        TEAL,
    )
    add_section_label(slide, "覆盖的测试入口", Inches(7.0), Inches(1.55), COPPER)
    add_bullets(
        slide,
        [
            "test_start_web_for_vision_result：启动原 Web 端。",
            "test_color_detect_result_to_web_with_base64：颜色结果和图像发送到前端。",
            "test_qr_detect_result_to_web_with_base64：QR 结果和最终彩色帧发送到前端。",
            "test_cross_detect_result_to_web_with_base64：当前发送占位结果，保证链路可测。",
        ],
        Inches(7.05),
        Inches(1.95),
        Inches(5.25),
        Inches(2.1),
        11.2,
        INK,
        gap=0.07,
        accent=GREEN,
    )
    add_panel(slide, Inches(7.05), Inches(4.68), Inches(5.18), Inches(0.72), fill=WHITE)
    add_textbox(
        slide,
        "注意：前端清空本地视图后不应自动全量同步历史，这是后续 Web 状态策略需要继续收敛的问题。",
        Inches(7.28),
        Inches(4.91),
        Inches(4.7),
        Inches(0.24),
        9.8,
        RED,
        True,
    )
    add_footer(slide, 8)


def slide_uart(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_title(slide, "串口通信测试：虚拟串口验证收发闭环", "使用 socat 创建虚拟串口对，验证命令触发与识别结果写回", dark=True)
    add_code_panel(
        slide,
        "虚拟串口启动命令",
        [
            "rm -f /tmp/rubovision_uart_a /tmp/rubovision_uart_b",
            "socat -d -d \\",
            "  pty,raw,echo=0,link=/tmp/rubovision_uart_a \\",
            "  pty,raw,echo=0,link=/tmp/rubovision_uart_b",
            "",
            "发送命令：printf 'a1\\n' > /tmp/rubovision_uart_b",
            "监听回写：cat /tmp/rubovision_uart_b",
        ],
        Inches(0.82),
        Inches(1.48),
        Inches(5.85),
        Inches(4.9),
        COPPER,
    )
    add_section_label(slide, "本次验证结论", Inches(7.05), Inches(1.58), TEAL)
    add_bullets(
        slide,
        [
            "配置指到 /tmp/rubovision_uart_a 后，a1 能触发 color_detect。",
            "关闭 GPIO 后，颜色识别完成并向串口对端写回 unknown。",
            "原配置 /dev/ttyV0 在当前机器不存在，因此主程序启动时报 open uart 失败。",
            "在非树莓派环境 returns.gpio=true 会先失败于 GPIO 初始化，导致摄像头尚未打开。",
        ],
        Inches(7.12),
        Inches(1.98),
        Inches(5.0),
        Inches(2.1),
        11.4,
        RGBColor(235, 240, 248),
        gap=0.08,
        accent=TEAL,
    )
    add_panel(slide, Inches(7.05), Inches(5.05), Inches(5.2), Inches(0.72), fill=rgb("111827"), line=rgb("3B4758"))
    add_textbox(slide, "定位结果：串口链路可用；摄像头未启用的根因是 GPIO/串口配置环境不匹配。", Inches(7.3), Inches(5.29), Inches(4.75), Inches(0.22), 10.2, WHITE, True)
    add_footer(slide, 9, dark=True)


def slide_issues(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_2)
    add_title(slide, "问题复盘与当前边界", "把已经暴露的问题纳入下一阶段计划，而不是用临时封装掩盖")
    widths = [Inches(2.25), Inches(4.25), Inches(3.1), Inches(2.65)]
    x = Inches(0.62)
    y = Inches(1.45)
    add_table_header(slide, ["问题", "现象", "定位/结论", "处理方向"], x, y, widths, fill=INK)
    rows = [
        ("Web 历史同步", "前端清空本地视图后收到新消息会重新带出历史", "前端本地状态和后端 history 快照语义需区分", "收敛同步策略"),
        ("串口未触发摄像头", "/dev/ttyV0 不存在时 UartSource 启动失败", "配置环境不匹配，命令未进入任务链路", "部署前确认串口路径"),
        ("GPIO 阻断识别", "非树莓派环境报 Unknown Raspberry Pi model", "GPIO 初始化发生在摄像头打开前", "测试时关闭 gpio 或加环境策略"),
        ("Vision 维护成本", "旧单文件过大，找测试与找功能困难", "拆分后职责清晰，仍需控制辅助函数数量", "继续按功能归档"),
    ]
    y += Inches(0.38)
    for i, row in enumerate(rows):
        add_table_row(slide, row, x, y, widths, row_h=0.72, fill=WHITE if i % 2 == 0 else rgb("F1F4F8"))
        y += Inches(0.72)
    add_section_label(slide, "阶段约束", Inches(0.82), Inches(5.4), COPPER)
    add_bullets(
        slide,
        [
            "只删除新增的冗余实现，不删除项目原有架构占位。",
            "优先用配置和测试暴露真实环境问题，避免为了“看起来通用”增加复杂抽象。",
            "cross 功能暂不实现算法，先保留可编译、可调度、可展示的函数体。",
        ],
        Inches(0.92),
        Inches(5.76),
        Inches(11.4),
        Inches(0.9),
        10.8,
        INK,
        gap=0.02,
        accent=RED,
    )
    add_footer(slide, 10)


def slide_verification(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_title(slide, "验证记录：已完成与待现场确认", "中期阶段已覆盖编译、单测、Web POST、虚拟串口；真实硬件仍需现场复测", dark=True)
    checks = [
        ("格式检查", "cargo fmt --check", "通过", TEAL),
        ("单元测试", "cargo test", "通过（GUI/硬件测试 ignored）", GREEN),
        ("Web 推送", "POST /message", "HTTP 202，可接收 base64 图像", BLUE),
        ("虚拟串口", "socat + a1/b2", "命令分发成功，结果可写回", COPPER),
        ("真实 GPIO", "rppal::Gpio", "当前 PC 环境不可用", RED),
        ("真实矿机", "摄像头/UART/GPIO 联动", "下一阶段现场确认", RGBColor(150, 160, 178)),
    ]
    for i, (name, cmd, status, color) in enumerate(checks):
        x = Inches(0.82) + Inches(4.18) * (i % 3)
        y = Inches(1.55) + Inches(1.58) * (i // 3)
        add_panel(slide, x, y, Inches(3.75), Inches(1.05), fill=rgb("232936"), line=rgb("3A4657"))
        add_solid_block(slide, color, x, y, Inches(3.75), Inches(0.08))
        add_textbox(slide, name, x + Inches(0.2), y + Inches(0.24), Inches(1.2), Inches(0.22), 12, WHITE, True)
        add_textbox(slide, status, x + Inches(1.48), y + Inches(0.24), Inches(2.05), Inches(0.22), 9.8, color, True, PP_ALIGN.RIGHT)
        add_textbox(slide, cmd, x + Inches(0.2), y + Inches(0.64), Inches(3.28), Inches(0.18), 8.7, RGBColor(210, 218, 230))
    add_code_panel(
        slide,
        "建议现场运行顺序",
        [
            "1. 确认 device.toml: serial=/dev/ttyV0 或实际串口",
            "2. 确认 /dev/video2 /dev/video4 与摄像头对应关系",
            "3. 在树莓派/矿机上启用 returns.gpio=true",
            "4. 发送 a1 / b2，核对 Web、UART 回写和 GPIO 状态灯",
        ],
        Inches(1.0),
        Inches(5.15),
        Inches(11.3),
        Inches(1.0),
        TEAL,
    )
    add_footer(slide, 11, dark=True)


def slide_next_plan(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_image(slide, ASSET_DIR / "mining-field-test.png")
    add_overlay(slide, BG, transparency=22)
    add_solid_block(slide, rgb("0F172A"), 0, 0, Inches(6.4), SLIDE_H)
    add_solid_block(slide, TEAL, Inches(6.4), 0, Inches(0.035), SLIDE_H)
    add_title(slide, "下一阶段计划", "从开发环境验证推进到真实矿机联调与稳定版本发布", dark=True)
    phases = [
        ("1", "现场环境校准", "确认串口、摄像头、GPIO、光照与安装角度。"),
        ("2", "前端状态收敛", "修正清空本地视图后自动同步历史的问题。"),
        ("3", "识别效果优化", "调整颜色阈值、ROI、QR 距离和连续识别策略。"),
        ("4", "cross 功能落地", "实现黑色轮廓/路口识别，并补齐单功能测试。"),
        ("5", "稳定版本整理", "冻结配置样例，补 README、测试记录和演示材料。"),
    ]
    y = Inches(1.7)
    for idx, title, desc in phases:
        add_chip(slide, idx, Inches(0.9), y, Inches(0.36), TEAL)
        add_textbox(slide, title, Inches(1.45), y - Inches(0.01), Inches(2.4), Inches(0.22), 13, WHITE, True)
        add_textbox(slide, desc, Inches(1.45), y + Inches(0.32), Inches(4.45), Inches(0.22), 10, RGBColor(222, 230, 242))
        y += Inches(0.86)
    add_panel(slide, Inches(0.92), Inches(6.18), Inches(5.25), Inches(0.58), fill=rgb("111827"), line=rgb("3B4758"))
    add_textbox(slide, "目标交付：可现场复现、可调试、可继续拓展的 RuboVision Engine 稳定版本。", Inches(1.12), Inches(6.38), Inches(4.85), Inches(0.18), 9.8, WHITE, True)
    add_footer(slide, 12, dark=True)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    for builder in [
        slide_cover,
        slide_context,
        slide_architecture,
        slide_returns_config,
        slide_vision_refactor,
        slide_vision_functions,
        slide_test_design,
        slide_web_debug,
        slide_uart,
        slide_issues,
        slide_verification,
        slide_next_plan,
    ]:
        builder(prs)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
