from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSET_DIR = ROOT / "assets"
OUT = ROOT / "RuboVision-Engine-midterm-report.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = "Microsoft YaHei"
FONT_FALLBACK = "Noto Sans CJK SC"

BG = RGBColor(18, 23, 34)
BG_2 = RGBColor(245, 247, 250)
INK = RGBColor(28, 33, 43)
MUTED = RGBColor(101, 111, 128)
WHITE = RGBColor(255, 255, 255)
COPPER = RGBColor(226, 145, 55)
TEAL = RGBColor(59, 184, 198)
GREEN = RGBColor(86, 179, 123)
BLUE = RGBColor(92, 131, 219)
RED = RGBColor(205, 92, 92)
LINE = RGBColor(218, 224, 233)


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_full_image(slide, image_path: Path):
    slide.shapes.add_picture(str(image_path), 0, 0, width=SLIDE_W, height=SLIDE_H)


def add_overlay(slide, color=BG, transparency=18, x=0, y=0, w=SLIDE_W, h=SLIDE_H):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.fill.background()
    return shape


def add_solid_block(slide, color, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    text,
    x,
    y,
    w,
    h,
    font_size=24,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing=1.05,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None, dark=False):
    color = WHITE if dark else INK
    add_textbox(slide, title, Inches(0.72), Inches(0.42), Inches(8.6), Inches(0.52), 25, color, True)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.74), Inches(0.94), Inches(9.7), Inches(0.35), 11, WHITE if dark else MUTED)


def add_footer(slide, idx, dark=False):
    color = RGBColor(205, 211, 221) if dark else RGBColor(125, 135, 150)
    add_textbox(
        slide,
        f"RuboVision Engine | 中期报告 | {idx:02d}",
        Inches(10.65),
        Inches(7.13),
        Inches(2.0),
        Inches(0.18),
        8,
        color,
        align=PP_ALIGN.RIGHT,
    )


def add_chip(slide, text, x, y, w, color, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = text_color
    return shape


def add_bullets(slide, items, x, y, w, h, font_size=15, color=INK, gap=0.18, accent=COPPER):
    top = y
    row_h = Inches(0.42)
    for item in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, top + Inches(0.12), Inches(0.09), Inches(0.09))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()
        add_textbox(slide, item, x + Inches(0.22), top, w - Inches(0.22), row_h, font_size, color)
        top += row_h + Inches(gap)


def add_panel(slide, x, y, w, h, fill=WHITE, line=LINE, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def add_metric(slide, label, value, note, x, y, w, accent):
    add_panel(slide, x, y, w, Inches(1.32), fill=WHITE)
    add_textbox(slide, value, x + Inches(0.18), y + Inches(0.18), w - Inches(0.36), Inches(0.38), 25, accent, True)
    add_textbox(slide, label, x + Inches(0.18), y + Inches(0.63), w - Inches(0.36), Inches(0.26), 12, INK, True)
    add_textbox(slide, note, x + Inches(0.18), y + Inches(0.93), w - Inches(0.36), Inches(0.25), 9, MUTED)


def add_connector(slide, x1, y1, x2, y2, color=TEAL):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2.0)
    line.line.end_arrowhead = True
    return line


def add_node(slide, text, x, y, w, h, fill, color=WHITE, font_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(255, 255, 255)
    shape.line.transparency = 55
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = color
    return shape


def add_section_label(slide, text, x, y, color=COPPER):
    add_textbox(slide, text, x, y, Inches(3.0), Inches(0.25), 9, color, True)


def add_table_header(slide, labels, x, y, widths, fill=INK):
    left = x
    for label, width in zip(labels, widths):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, width, Inches(0.38))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill
        add_textbox(slide, label, left + Inches(0.08), y + Inches(0.08), width - Inches(0.16), Inches(0.16), 9, WHITE, True)
        left += width


def add_table_row(slide, values, x, y, widths, row_h=0.58, fill=WHITE):
    left = x
    for value, width in zip(values, widths):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, width, Inches(row_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = LINE
        shape.line.width = Pt(0.6)
        add_textbox(slide, value, left + Inches(0.08), y + Inches(0.08), width - Inches(0.16), Inches(row_h - 0.12), 9.3, INK)
        left += width


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_image(slide, ASSET_DIR / "cover-vision-engine.png")
    add_solid_block(slide, BG, 0, 0, Inches(5.85), SLIDE_H)
    add_solid_block(slide, rgb("233044"), Inches(5.85), 0, Inches(0.03), SLIDE_H)
    add_textbox(slide, "RuboVision Engine", Inches(0.72), Inches(1.15), Inches(5.7), Inches(0.55), 29, WHITE, True)
    add_textbox(slide, "中期报告", Inches(0.74), Inches(1.82), Inches(4.0), Inches(0.45), 24, COPPER, True)
    add_textbox(
        slide,
        "框架结构完善 · 功能拓展 · 矿机落地测试规划",
        Inches(0.76),
        Inches(2.48),
        Inches(5.8),
        Inches(0.42),
        14,
        RGBColor(230, 235, 244),
    )
    add_bullets(
        slide,
        ["项目组成员：2 人", "当前版本：Rust + OpenCV + Axum 的配置驱动视觉任务框架", "报告时间：2026 年 5 月"],
        Inches(0.78),
        Inches(4.7),
        Inches(5.3),
        Inches(1.3),
        12,
        RGBColor(232, 237, 246),
        gap=0.05,
        accent=TEAL,
    )
    add_footer(slide, 1, dark=True)


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_2)
    add_title(slide, "项目定位与建设目标", "从单点视觉程序升级为可配置、可扩展、可落地测试的矿机视觉任务框架")
    add_metric(slide, "核心链路", "Source→WebMessage", "事件进入、任务调度、执行与 Web 输出闭环", Inches(0.82), Inches(1.55), Inches(3.7), TEAL)
    add_metric(slide, "功能迁移", "2+1", "颜色识别、二维码识别已接入；路口识别保留扩展点", Inches(4.82), Inches(1.55), Inches(3.7), GREEN)
    add_metric(slide, "配置入口", "4 类", "bindings / device / func_param / web 分层配置", Inches(8.82), Inches(1.55), Inches(3.7), COPPER)
    add_section_label(slide, "项目建设目标", Inches(0.82), Inches(3.38))
    add_bullets(
        slide,
        [
            "将旧版车载视觉逻辑抽象为“来源、设备、函数、结果”的通用框架。",
            "通过配置文件决定触发源、设备实例和功能参数，减少主循环硬编码。",
            "面向矿机落地场景，打通 UART 指令、摄像头识别、GPIO 状态灯与 Web 调试。",
            "形成可复用框架：后续新增功能只需要注册设备、函数与绑定关系。",
        ],
        Inches(0.88),
        Inches(3.74),
        Inches(11.4),
        Inches(2.1),
        14,
    )
    add_footer(slide, 2)


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_title(slide, "中期完成情况", "目前已从旧版功能同步到新版架构，并补齐运行、回执和调试链路", dark=True)
    items = [
        ("框架链路", "完成 Source、TaskListener、TaskDispatcher、TaskExecutor、WebMessage 串联。", TEAL),
        ("配置驱动", "bindings.toml 映射任务；device.toml 管理 UART/Camera；func_param.toml 管理函数参数。", COPPER),
        ("视觉功能", "颜色识别完成 HSV 动态配置、ROI、面积比例、连续稳定计数；二维码识别完成灰度预处理与 quircs 解码。", GREEN),
        ("硬件回执", "统一串口配置用于监听与识别结果回写；GPIO 状态灯已接入任务执行过程。", BLUE),
        ("Web 调试", "统一 WebMessage 输出，支持 latest/history，消息持久化到 jsonl，Web 关闭时不阻塞执行端。", RED),
    ]
    x = Inches(0.85)
    y = Inches(1.55)
    for title, desc, color in items:
        add_panel(slide, x, y, Inches(11.7), Inches(0.82), fill=rgb("232936"), line=rgb("374154"))
        add_chip(slide, title, x + Inches(0.18), y + Inches(0.22), Inches(1.25), color)
        add_textbox(slide, desc, x + Inches(1.66), y + Inches(0.19), Inches(9.6), Inches(0.38), 12.2, RGBColor(236, 240, 248))
        y += Inches(0.96)
    add_footer(slide, 3, dark=True)


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_image(slide, ASSET_DIR / "framework-architecture.png")
    add_solid_block(slide, BG, 0, 0, SLIDE_W, Inches(1.35))
    add_title(slide, "框架结构：从事件到结果的消息链路", "新版采用统一任务链路，功能不再写死在主循环中", dark=True)
    y = Inches(2.15)
    xs = [Inches(0.75), Inches(3.1), Inches(5.45), Inches(7.8), Inches(10.15)]
    labels = ["Source\nUART/Timer/Web", "TaskListener\n接收事件", "Dispatcher\n匹配设备+函数", "Executor\n阻塞任务隔离", "WebMessage\n结果输出"]
    colors = [TEAL, BLUE, COPPER, GREEN, RED]
    for i, (x, label, color) in enumerate(zip(xs, labels, colors)):
        add_node(slide, label, x, y, Inches(1.72), Inches(0.95), color, WHITE, 12)
        if i < len(xs) - 1:
            add_connector(slide, x + Inches(1.72), y + Inches(0.48), xs[i + 1], y + Inches(0.48), RGBColor(198, 219, 232))
    add_panel(slide, Inches(0.85), Inches(4.15), Inches(11.55), Inches(1.55), fill=rgb("121722"), line=RGBColor(105, 119, 139), transparency=0)
    add_bullets(
        slide,
        [
            "事件类型：UsualEvent / DebugEvent / OtherEvent，为后续来源扩展预留空间。",
            "调度依据：bindings 中的 task_id、source_key、device_id、function_id。",
            "执行方式：视觉识别等阻塞任务使用 spawn_blocking，降低对异步事件循环的影响。",
        ],
        Inches(1.2),
        Inches(4.42),
        Inches(10.8),
        Inches(1.1),
        12,
        RGBColor(234, 239, 247),
        gap=0.02,
        accent=COPPER,
    )
    add_footer(slide, 4, dark=True)


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_2)
    add_title(slide, "功能拓展：视觉识别与硬件回执", "中期阶段重点完成旧版核心能力迁移，并把输出纳入统一框架")
    cards = [
        ("颜色识别", ["OpenCV 摄像头读取", "圆形 ROI 与 HSV 筛选", "颜色数量由配置动态注册", "稳定计数与面积比例过滤"], TEAL),
        ("二维码识别", ["灰度预处理", "quircs 解码", "任务编号解析", "串口结果回写"], GREEN),
        ("硬件联动", ["UART 命令触发", "统一串口配置", "GPIO 运行/任务状态灯", "异常记录到 WebMessage"], COPPER),
        ("待扩展功能", ["cross_detect 已注册", "当前返回占位值 0", "后续接入黑色轮廓/路口识别", "作为框架扩展样例"], BLUE),
    ]
    x0 = Inches(0.8)
    y0 = Inches(1.55)
    for i, (title, bullets, color) in enumerate(cards):
        x = x0 + Inches(3.05) * i
        add_panel(slide, x, y0, Inches(2.75), Inches(4.75), fill=WHITE)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y0, Inches(2.75), Inches(0.12))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        add_textbox(slide, title, x + Inches(0.18), y0 + Inches(0.34), Inches(2.25), Inches(0.28), 15, INK, True)
        add_bullets(slide, bullets, x + Inches(0.23), y0 + Inches(0.98), Inches(2.25), Inches(2.9), 10.6, INK, gap=0.09, accent=color)
    add_footer(slide, 5)


def slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_2)
    add_title(slide, "配置化扩展：让新增功能变成可复用流程", "当前配置文件已经拆分出触发、设备、函数参数与 Web 调试入口")
    left_x = Inches(0.82)
    top = Inches(1.58)
    config_items = [
        ("bindings.toml", "将 source_key 映射到 task_id + device_id + function_id"),
        ("device.toml", "声明进程级唯一 UART 与多个 Camera 设备实例"),
        ("func_param.toml", "注册功能函数、返回通道、HSV/ROI/GPIO 等参数"),
        ("web.yaml", "控制 Web 调试面板 host、port 与启用状态"),
    ]
    for name, desc in config_items:
        add_panel(slide, left_x, top, Inches(5.45), Inches(0.72), fill=WHITE)
        add_textbox(slide, name, left_x + Inches(0.2), top + Inches(0.14), Inches(1.6), Inches(0.22), 12, COPPER, True)
        add_textbox(slide, desc, left_x + Inches(1.95), top + Inches(0.14), Inches(3.2), Inches(0.26), 10.2, INK)
        top += Inches(0.88)
    add_section_label(slide, "新增能力的推荐路径", Inches(7.0), Inches(1.58))
    steps = [
        ("1", "定义设备类型", "在 device 侧注册硬件实例与必要参数校验"),
        ("2", "实现功能函数", "在 func 侧封装识别逻辑，返回 WebMessage"),
        ("3", "绑定触发来源", "在 bindings 中配置指令、设备与函数关系"),
        ("4", "补齐样例与测试", "提供示例配置、Web 调试、串口/GPIO 验证"),
    ]
    y = Inches(1.95)
    for idx, title, desc in steps:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.05), y, Inches(0.38), Inches(0.38))
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()
        add_textbox(slide, idx, Inches(7.05), y + Inches(0.06), Inches(0.38), Inches(0.13), 10, WHITE, True, PP_ALIGN.CENTER)
        add_textbox(slide, title, Inches(7.58), y - Inches(0.02), Inches(2.3), Inches(0.25), 12, INK, True)
        add_textbox(slide, desc, Inches(7.58), y + Inches(0.28), Inches(4.2), Inches(0.25), 9.5, MUTED)
        if idx != "4":
            add_connector(slide, Inches(7.24), y + Inches(0.42), Inches(7.24), y + Inches(0.78), LINE)
        y += Inches(0.9)
    add_footer(slide, 6)


def slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_image(slide, ASSET_DIR / "mining-field-test.png")
    add_solid_block(slide, BG, 0, 0, Inches(6.35), SLIDE_H)
    add_solid_block(slide, rgb("233044"), Inches(6.35), 0, Inches(0.03), SLIDE_H)
    add_title(slide, "矿机落地测试方案", "把框架从开发环境推进到真实设备链路验证", dark=True)
    add_section_label(slide, "测试重点", Inches(0.82), Inches(1.68), TEAL)
    add_bullets(
        slide,
        [
            "UART 指令触发：验证 a1 / b2 等命令能稳定触发对应任务。",
            "摄像头识别：在矿机安装视角下校准 ROI、光照、颜色阈值与二维码距离。",
            "回执链路：验证识别结果串口回写、GPIO 状态灯和 Web 历史消息一致。",
            "异常场景：断摄像头、串口异常、无目标、误识别、连续触发等稳定性测试。",
        ],
        Inches(0.9),
        Inches(2.05),
        Inches(5.35),
        Inches(2.2),
        12.4,
        RGBColor(236, 240, 248),
        gap=0.08,
        accent=COPPER,
    )
    add_panel(slide, Inches(0.9), Inches(5.45), Inches(5.28), Inches(0.72), fill=rgb("111722"), line=RGBColor(116, 132, 153), transparency=0)
    add_textbox(slide, "目标：形成可复现的矿机联调记录，为稳定版本和示例工程提供依据。", Inches(1.1), Inches(5.68), Inches(4.9), Inches(0.2), 10.5, WHITE, True)
    add_footer(slide, 7, dark=True)


def slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_2)
    add_title(slide, "未来规划：第 12-16 周四阶段推进", "围绕矿机落地测试、功能完善、完整框架、示例与稳定版本发布")
    phases = [
        ("第12周", "联调准备", "整理硬件清单、确认串口/摄像头路径、制定矿机场景测试表。", TEAL, 0, 1),
        ("第13周", "落地测试", "完成矿机现场颜色/二维码识别测试，记录误识别与延迟问题。", GREEN, 1, 2),
        ("第14-15周", "功能完善", "优化阈值、异常处理、Web 调试、cross_detect 与框架扩展示例。", COPPER, 2, 4),
        ("第16周", "稳定发布", "冻结配置样例，补齐 README/示例/测试记录，发布稳定版本。", BLUE, 4, 5),
    ]
    x0 = Inches(1.0)
    y0 = Inches(2.15)
    track_w = Inches(10.95)
    week_w = track_w / 5
    for i, week in enumerate(["12", "13", "14", "15", "16"]):
        add_textbox(slide, f"W{week}", x0 + week_w * i, Inches(1.63), week_w, Inches(0.22), 10, MUTED, True, PP_ALIGN.CENTER)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0 + week_w * i + week_w / 2, Inches(1.94), Inches(0.01), Inches(3.4))
        line.fill.solid()
        line.fill.fore_color.rgb = LINE
        line.line.fill.background()
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, Inches(2.0), track_w, Inches(0.04))
    base.fill.solid()
    base.fill.fore_color.rgb = LINE
    base.line.fill.background()
    y = y0
    for week_label, title, desc, color, start, end in phases:
        add_textbox(slide, week_label, Inches(0.78), y + Inches(0.12), Inches(1.0), Inches(0.23), 10, color, True)
        bar_x = x0 + week_w * start + Inches(0.05)
        bar_w = week_w * (end - start) - Inches(0.1)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, y, bar_w, Inches(0.52))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(slide, title, bar_x + Inches(0.14), y + Inches(0.15), bar_w - Inches(0.28), Inches(0.18), 10, WHITE, True, PP_ALIGN.CENTER)
        add_textbox(slide, desc, Inches(1.05), y + Inches(0.68), Inches(10.9), Inches(0.22), 9.3, INK)
        y += Inches(0.98)
    add_footer(slide, 8)


def slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_2)
    add_title(slide, "双人成员分工与风险控制", "两人协同推进：一人偏框架与发布，一人偏视觉与落地测试")
    widths = [Inches(1.45), Inches(3.8), Inches(3.8), Inches(3.0)]
    x = Inches(0.75)
    y = Inches(1.55)
    add_table_header(slide, ["阶段", "成员 A：框架/工程化", "成员 B：视觉/硬件联调", "共同输出"], x, y, widths)
    rows = [
        ("第12周", "整理架构文档、测试表与示例配置模板", "核对矿机摄像头/UART/GPIO 接线与参数", "联调准备清单"),
        ("第13周", "记录 WebMessage、日志、串口回执问题", "完成现场识别测试与阈值初调", "矿机测试记录"),
        ("第14-15周", "完善错误处理、示例工程、Web 调试体验", "完善 cross_detect/识别鲁棒性与硬件适配", "功能完善版本"),
        ("第16周", "整理 README、版本说明、发布包", "复测关键用例、输出演示素材", "稳定版本与答辩材料"),
    ]
    y += Inches(0.38)
    for i, row in enumerate(rows):
        add_table_row(slide, row, x, y, widths, row_h=0.66, fill=WHITE if i % 2 == 0 else rgb("F1F4F8"))
        y += Inches(0.66)
    add_section_label(slide, "主要风险与应对", Inches(0.82), Inches(5.18), COPPER)
    add_bullets(
        slide,
        [
            "现场光照、震动和安装角度导致识别波动：通过 ROI、阈值配置、连续稳定计数和测试样本沉淀解决。",
            "硬件环境差异导致串口/摄像头路径变化：保留配置化入口，稳定版本提供多套示例配置。",
            "功能扩展后框架边界变模糊：用 cross_detect 作为扩展样例，明确新增功能的注册流程。",
        ],
        Inches(0.95),
        Inches(5.55),
        Inches(11.5),
        Inches(1.25),
        10.8,
        INK,
        gap=0.02,
        accent=RED,
    )
    add_footer(slide, 9)


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_title(slide, "预期交付：完整框架与稳定版本", "第 16 周形成可展示、可复现、可继续拓展的阶段成果", dark=True)
    deliverables = [
        ("完整框架", "Source / Device / Func / WebMessage 链路稳定，新增功能流程清晰。", TEAL),
        ("示例配置", "提供颜色识别、二维码识别、矿机测试、调试模式等示例。", COPPER),
        ("稳定版本", "冻结核心接口，补齐测试记录、异常处理和版本说明。", GREEN),
        ("演示材料", "形成矿机落地测试记录、Web 调试截图、答辩演示脚本。", BLUE),
    ]
    x = Inches(0.85)
    for i, (title, desc, color) in enumerate(deliverables):
        y = Inches(1.65) + Inches(1.07) * i
        add_panel(slide, x, y, Inches(11.75), Inches(0.78), fill=rgb("232936"), line=rgb("3A4657"))
        add_chip(slide, title, x + Inches(0.2), y + Inches(0.23), Inches(1.3), color)
        add_textbox(slide, desc, x + Inches(1.78), y + Inches(0.22), Inches(9.7), Inches(0.24), 12, RGBColor(238, 242, 248))
    add_textbox(
        slide,
        "最终目标：把 RuboVision Engine 从“能跑的功能集合”推进为“有边界、有示例、有发布物的视觉任务框架”。",
        Inches(1.05),
        Inches(6.34),
        Inches(11.2),
        Inches(0.38),
        14,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_footer(slide, 10, dark=True)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    # Remove default slide if present in some environments.
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    for builder in [
        slide_1,
        slide_2,
        slide_3,
        slide_4,
        slide_5,
        slide_6,
        slide_7,
        slide_8,
        slide_9,
        slide_10,
    ]:
        builder(prs)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
